import os
import re
import hashlib
from flask import Flask, request, jsonify
from flask_cors import CORS
import fitz  # PyMuPDF for PDF processing
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import CountVectorizer
import numpy as np
import requests

app = Flask(__name__)
CORS(app)
TEMP_DIR = "temp_files"
os.makedirs(TEMP_DIR, exist_ok=True)

RUBRIC_WEIGHTS = {
    "originality": 0.4,
    "structure": 0.2,
    "grammar": 0.2,
    "completeness": 0.2
}

def preprocess_text(text):
    """Lowercase and remove punctuation."""
    if not text:
        return ""
    text = text.lower()
    text = re.sub(r"[\W_]+", " ", text)
    return text

def extract_text_from_pdf(file):
    """Extract text from PDF using PyMuPDF."""
    text = ""
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page in pdf_document:
        text += page.get_text()
    return text.strip()

def extract_text_from_docx(file):
    """Extract text from DOCX files."""
    doc = Document(file)
    return " ".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    """Extract text from PPTX files."""
    presentation = Presentation(file)
    return " ".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

def process_file(file, filename):
    """Process supported file types."""
    if filename.endswith(".pdf"):
        return extract_text_from_pdf(file)
    elif filename.endswith(".docx"):
        return extract_text_from_docx(file)
    elif filename.endswith(".pptx"):
        return extract_text_from_pptx(file)
    return None

def cosine_similarity(text1, text2):
    """Calculate the cosine similarity between two texts."""
    vectorizer = CountVectorizer()
    vectors = vectorizer.fit_transform([text1, text2]).toarray()
    dot_product = np.dot(vectors[0], vectors[1])
    magnitude1 = np.linalg.norm(vectors[0])
    magnitude2 = np.linalg.norm(vectors[1])
    if magnitude1 == 0 or magnitude2 == 0:
        return 0
    cosine_sim = dot_product / (magnitude1 * magnitude2)
    return cosine_sim

def jaccard_similarity(text1, text2):
    """Calculate Jaccard similarity between two texts."""
    set1 = set(text1.split())
    set2 = set(text2.split())
    intersection = len(set1.intersection(set2))
    union = len(set1.union(set2))
    return intersection / union if union != 0 else 0

def evaluate_rubric(text):
    """Evaluate text based on rubric criteria."""
    words = text.split()
    num_words = len(words)
    
    originality = 1.0  # Placeholder for more advanced originality detection
    structure = min(1.0, num_words / 500)  # Simple structure heuristic (adjust as needed)
    grammar = 1.0 if num_words > 20 else 0.5  # Placeholder for grammar checking
    completeness = min(1.0, num_words / 1000)  # Adjust based on expected word count

    rubric_scores = {
        "originality": originality,
        "structure": structure,
        "grammar": grammar,
        "completeness": completeness
    }

    final_score = sum(rubric_scores[criterion] * weight for criterion, weight in RUBRIC_WEIGHTS.items()) * 100

    return {**rubric_scores, "final_score": round(final_score, 2)}

@app.route('/upload', methods=['POST'])
def upload():
    try:
        file_details = request.json.get('fileDetails', [])
        print("Uploading file details", file_details)
        if 'fileDetails' not in request.json:
            return jsonify({"success": False, "message": "Invalid request data"}), 400

        file_paths = []
        files = {}
        for file in file_details:
            student_id = file.get('studentId')
            file_url = file.get('fileUrl')

            if not file_url or not student_id:
                continue

            response = requests.get(file_url)
            if response.status_code == 200:
                file_name = f"{os.path.basename(file_url)}"
                file_path = os.path.join(TEMP_DIR, file_name)
                files[file_name] = student_id
                with open(file_path, 'wb') as f:
                    f.write(response.content)
                file_paths.append(file_path)

            else:
                print(f"Failed to download file from: {file_url}")
    
        uploaded_files = []
        assignments = {}
        rubric_scores = {}
        for file in file_paths:
            uploaded_files.append(file)

        for file in uploaded_files:
            with open(file, "rb") as f:
                text = process_file(f, file)
                if text:
                    clean_text = preprocess_text(text)
                    assignments[os.path.basename(file)] = clean_text
                    rubric_scores[os.path.basename(file)] = evaluate_rubric(clean_text)
                else:
                    text = process_file(file, file)
        
        # Calculate similarities
        result = []
        assignment_names = list(assignments.keys())
        print(assignment_names)
        for i in range(len(assignment_names)):
            for j in range(i + 1, len(assignment_names)):
                name1 = assignment_names[i]
                name2 = assignment_names[j]
                print(name1, name2)
                # Calculate Cosine Similarity
                cosine_sim = cosine_similarity(assignments[name1], assignments[name2])
                # Calculate Jaccard Similarity
                jaccard_sim = jaccard_similarity(assignments[name1], assignments[name2])

                # Combine both similarities into a final combined score
                combined_sim = (cosine_sim + jaccard_sim) / 2

                result.append({
                    "studentId1": files[name1],
                    "studentId2": files[name2],
                    "Assignment1": name1,
                    "Assignment2": name2,
                    "Cosine Similarity (%)": round(float(cosine_sim) * 100, 2),
                    "Jaccard Similarity (%)": round(float(jaccard_sim) * 100, 2),
                    "Combined Similarity (%)": round(float(combined_sim) * 100, 2)
                })
        
        # Include rubric scores in response
        print(result)
        rubric_results = []
        for assignment, scores in rubric_scores.items():
            rubric_results.append({
                "Assignment": assignment,
                "StudentId": files[assignment],
                "Originality Score": round(scores["originality"] * 100, 2),
                "Structure Score": round(scores["structure"] * 100, 2),
                "Grammar Score": round(scores["grammar"] * 100, 2),
                "Completeness Score": round(scores["completeness"] * 100, 2),
                "Final Rubric Score (%)": scores["final_score"]
            })
        print('hello', rubric_results)
        return jsonify({"success": True, "results": result, "rubricResults": rubric_results})
    
    except Exception as e:
        print(e)
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
